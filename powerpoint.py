from pptx import Presentation
from pptx.util import Inches
import os

#width and height of images I want to display in the powerpoint
width = Inches(2.5)
height = Inches(3.5)
listPNGS = []

#all of the lists of treatment groups image names
KRAS_NO_ARS = []
KRAS_ARS = []
HRAS_NO_ARS = []
HRAS_ARS = []
NRAS_NO_ARS = []
NRAS_ARS = []

#intializes each of the treatment group arrays and creates the final presentation
def run():
	for pic in listPNGS:
		if 'NRAS NO ARS' in pic:
			NRAS_NO_ARS.append(pic)
		elif 'NRAS ARS' in pic:
			NRAS_ARS.append(pic)
		elif 'HRAS ARS' in pic:
			HRAS_ARS.append(pic)
		elif 'HRAS NO ARS' in pic:
			HRAS_NO_ARS.append(pic)
		elif 'KRAS ARS' in pic:
			KRAS_ARS.append(pic)
		elif 'KRAS NO ARS' in pic:
			KRAS_NO_ARS.append(pic)

	minimum = findMin()

	prs = Presentation()

	for i in range(minimum):
		img_path_KRAS_NO_ARS = KRAS_NO_ARS[i]

		img_path_KRAS_ARS = KRAS_ARS[i]

		img_path_NRAS_NO_ARS = NRAS_NO_ARS[i]

		img_path_NRAS_ARS = NRAS_ARS[i]

		img_path_HRAS_NO_ARS = HRAS_NO_ARS[i]

		img_path_HRAS_ARS = HRAS_ARS[i]

		title_slide_layout = prs.slide_layouts[6]
		slide = prs.slides.add_slide(title_slide_layout)
		
		left = Inches(.5)
		top = Inches(.25)
		
		pic = slide.shapes.add_picture(img_path_KRAS_NO_ARS, left, top, width=width, height=height)
		
		left = Inches(3.5)
		pic = slide.shapes.add_picture(img_path_NRAS_NO_ARS, left, top, width=width, height=height)
		
		top = Inches(3.75)
		left = Inches(.5)
		
		pic = slide.shapes.add_picture(img_path_KRAS_ARS, left, top, width=width, height=height)
		
		left = Inches(3.5)
		
		pic = slide.shapes.add_picture(img_path_NRAS_ARS, left, top, width=width, height=height)

		top = Inches(.25)
		left = Inches(7)
		pic = slide.shapes.add_picture(img_path_HRAS_NO_ARS, left, top, width=width, height=height)

		top = Inches(3.75)
		left = Inches(7)
		pic = slide.shapes.add_picture(img_path_HRAS_ARS, left, top, width=width, height=height)

	prs.save('test.pptx')

#finds the minimum size of all of the treatment group arrays via some simple comparisons
def findMin():
	minimum = len(KRAS_NO_ARS)
	if len(KRAS_ARS) < minimum:
		minimum = len(KRAS_ARS)

	if len(HRAS_NO_ARS) < minimum:
		minimum = len(HRAS_NO_ARS)

	if len(HRAS_ARS) < minimum:
		minimum = len(HRAS_ARS)

	if len(NRAS_NO_ARS) < minimum:
		minimum = len(NRAS_NO_ARS)

	if len(NRAS_ARS) < minimum:
		minimum = len(NRAS_ARS)
	return minimum


if __name__ == "__main__":
	listPNGS = os.listdir()
	toRemove = []
	for element in listPNGS:
		if "png" not in element:
			toRemove.append(element)
	for picture in toRemove:
		listPNGS.remove(picture)
	run()